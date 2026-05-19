"""Extrae texto plano del PDF y PPTX de la presentación Trongkai.

Volca a docs/dumps/ para análisis del equipo y de los agentes.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
CTX = ROOT / "contexto"
OUT = ROOT / "docs" / "dumps"
OUT.mkdir(parents=True, exist_ok=True)

PDF = CTX / "Presentacion_Trongkai_2026-10.pdf"
PPTX = CTX / "Presentacion_Trongkai_2025-10-22.pptx"


def extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    lines = [f"# {path.name}", f"# {len(reader.pages)} páginas\n"]
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            lines.append(f"\n--- página {i} ---\n{text}")
    return "\n".join(lines)


def extract_pptx(path: Path) -> str:
    pres = Presentation(str(path))
    lines = [f"# {path.name}", f"# {len(pres.slides)} slides\n"]
    for i, slide in enumerate(pres.slides, 1):
        lines.append(f"\n--- slide {i} ---")
        # Notas
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[notas]: {notes}")
        # Texto de los shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    lines.append(txt)
            elif shape.shape_type == 19:  # TABLE
                tbl = shape.table
                for row in tbl.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        lines.append("\t".join(cells))
    return "\n".join(lines)


def main() -> int:
    if PDF.exists():
        out = OUT / "presentacion_pdf.txt"
        out.write_text(extract_pdf(PDF), encoding="utf-8")
        print(f"PDF ->{out} ({out.stat().st_size:,} bytes)")
    else:
        print(f"WARN: no se encuentra {PDF}", file=sys.stderr)

    if PPTX.exists():
        out = OUT / "presentacion_pptx.txt"
        out.write_text(extract_pptx(PPTX), encoding="utf-8")
        print(f"PPTX ->{out} ({out.stat().st_size:,} bytes)")
    else:
        print(f"WARN: no se encuentra {PPTX}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
